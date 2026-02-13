import os
import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


# ============================================================
# UTILITIES
# ============================================================

def format_months(months):
    formatted = []
    for m in months:
        try:
            formatted.append(pd.to_datetime(str(m)).strftime("%b-%y"))
        except:
            formatted.append(str(m))
    return formatted


def save_line_chart(data, path, title=None, secondary=None):

    if not data:
        return False

    months = sorted(data.keys())
    values = [data[m] for m in months]
    months_fmt = format_months(months)

    fig, ax1 = plt.subplots(figsize=(4, 3))

    ax1.plot(months_fmt, values, marker="o", linewidth=2)
    ax1.tick_params(axis="x", rotation=45, labelsize=7)
    ax1.tick_params(axis="y", labelsize=7)

    if secondary:
        ax2 = ax1.twinx()
        sec_values = [secondary.get(m, 0) for m in months]
        ax2.plot(months_fmt, sec_values, linestyle="--", linewidth=2)
        ax2.tick_params(axis="y", labelsize=7)

    if title:
        plt.title(title, fontsize=9)

    plt.tight_layout()
    plt.savefig(path, dpi=200)
    plt.close()
    return True


def save_combo_chart(combo_data, path, title=None):

    monthly_total = combo_data.get("monthly_total", {})
    region_pct = combo_data.get("monthly_region_pct", {})

    if not monthly_total:
        return False

    months = sorted(monthly_total.keys())
    months_fmt = format_months(months)

    regions = set()
    for m in region_pct.values():
        regions.update(m.keys())
    regions = sorted(regions)

    stacks = {r: [] for r in regions}
    totals = []

    for m in months:
        totals.append(monthly_total.get(m, 0))
        month_data = region_pct.get(m, {})
        for r in regions:
            stacks[r].append(month_data.get(r, 0))

    fig, ax1 = plt.subplots(figsize=(4, 3))

    bottom = [0] * len(months)
    for r in regions:
        ax1.bar(months_fmt, stacks[r], bottom=bottom, label=r)
        bottom = [i + j for i, j in zip(bottom, stacks[r])]

    ax1.set_ylim(0, 100)
    ax1.tick_params(axis="x", rotation=45, labelsize=6)
    ax1.tick_params(axis="y", labelsize=6)

    ax2 = ax1.twinx()
    ax2.plot(months_fmt, totals, color="black", marker="o", linewidth=2)
    ax2.tick_params(axis="y", labelsize=6)

    ax1.legend(fontsize=6)

    if title:
        plt.title(title, fontsize=9)

    plt.tight_layout()
    plt.savefig(path, dpi=200)
    plt.close()
    return True


def add_insights(slide, insights):

    if not insights:
        return

    box = slide.shapes.add_textbox(
        Inches(0.5),
        Inches(6.2),
        Inches(12),
        Inches(1.1)
    )

    tf = box.text_frame
    tf.word_wrap = True

    for i, line in enumerate(insights):
        if i == 0:
            tf.text = "• " + line
        else:
            p = tf.add_paragraph()
            p.text = "• " + line
            p.level = 0


# ============================================================
# MAIN PPT GENERATOR
# ============================================================

def generate_ppt(results):

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    os.makedirs("outputs", exist_ok=True)
    os.makedirs("temp", exist_ok=True)

    ap = results["ap_kpis"]
    ar = results["ar_kpis"]
    insights = results.get("insights", {})

    # ============================================================
    # SLIDE 1 – AP SPEND
    # ============================================================

    slide1 = prs.slides.add_slide(prs.slide_layouts[6])

    save_line_chart(ap["monthly_spend_trend"],
                    "temp/ap_s1_total.png",
                    "Monthly Spend",
                    secondary=ap["monthly_budget"])

    slide1.shapes.add_picture("temp/ap_s1_total.png",
                              Inches(0.5), Inches(0.8),
                              width=Inches(6), height=Inches(3))

    bus = ["Consumer", "Enterprise", "SMB"]

    positions = [(7, 0.8), (0.5, 3.8), (7, 3.8)]

    for i, bu in enumerate(bus):
        combo = ap["bu_region_combo"].get(bu, {})
        if combo:
            path = f"temp/ap_s1_{bu}.png"
            save_combo_chart(combo, path, f"{bu} BU Mix")
            slide1.shapes.add_picture(path,
                                      Inches(positions[i][0]),
                                      Inches(positions[i][1]),
                                      width=Inches(6),
                                      height=Inches(3))

    add_insights(slide1, insights.get("ap", {}).get("slide1", []))

    # ============================================================
    # SLIDE 2 – OPERATIONAL
    # ============================================================

    slide2 = prs.slides.add_slide(prs.slide_layouts[6])

    save_line_chart(ap["processing_time_trend"],
                    "temp/ap_s2_proc.png",
                    "Avg Processing Time")
    
    save_line_chart(ap["late_payment_pct"],
                    "temp/ap_s2_late.png",
                    "% Late Payment")
    
    save_line_chart(ap["dispute_pct"],
                    "temp/ap_s2_disp.png",
                    "%Dispute")
    

    slide2.shapes.add_picture("temp/ap_s2_proc.png",
                              Inches(0.5), Inches(0.8),
                              width=Inches(6), height=Inches(3))
    slide2.shapes.add_picture("temp/ap_s2_late.png",
                              Inches(7), Inches(0.8),
                              width=Inches(6), height=Inches(3))
    slide2.shapes.add_picture("temp/ap_s2_disp.png",
                              Inches(0.5), Inches(3.8),
                              width=Inches(6), height=Inches(3))
    


    add_insights(slide2, insights.get("ap", {}).get("slide2", []))

    # ============================================================
    # SLIDE 3 – COMPLIANCE
    # ============================================================

    slide3 = prs.slides.add_slide(prs.slide_layouts[6])

    save_line_chart(ap["maverick_pct"],
                    "temp/ap_s3_mav.png",
                    "Maverick %")

    slide3.shapes.add_picture("temp/ap_s3_mav.png",
                              Inches(0.5), Inches(1),
                              width=Inches(6), height=Inches(3))

    add_insights(slide3, insights.get("ap", {}).get("slide3", []))

    # ============================================================
    # SLIDE 4 – WORKING CAPITAL
    # ============================================================

    slide4 = prs.slides.add_slide(prs.slide_layouts[6])

    save_line_chart(ap["dpo_trend"],
                    "temp/ap_s4_dpo.png",
                    "DPO Trend",
                    secondary=ap["on_time_pct"])

    slide4.shapes.add_picture("temp/ap_s4_dpo.png",
                              Inches(0.5), Inches(1),
                              width=Inches(6), height=Inches(3))

    add_insights(slide4, insights.get("ap", {}).get("slide4", []))

    # ============================================================
    # SLIDE 5 – AR REVENUE & COLLECTION
    # ============================================================

    slide5 = prs.slides.add_slide(prs.slide_layouts[6])

    save_line_chart(ar["revenue_trend"],
                    "temp/ar_s1_rev.png",
                    "Revenue",
                    secondary=ar["collection_trend"])

    slide5.shapes.add_picture("temp/ar_s1_rev.png",
                              Inches(0.5), Inches(1),
                              width=Inches(6), height=Inches(3))

    add_insights(slide5, insights.get("ar", {}).get("slide1", []))

    # ============================================================
    # SLIDE 6 – AR SEGMENT MIX
    # ============================================================

    slide6 = prs.slides.add_slide(prs.slide_layouts[6])

    save_line_chart(ar["segment_mix_pct"],
                    "temp/ar_s2_seg.png",
                    "Segment Mix %")

    slide6.shapes.add_picture("temp/ar_s2_seg.png",
                              Inches(0.5), Inches(1),
                              width=Inches(6), height=Inches(3))

    add_insights(slide6, insights.get("ar", {}).get("slide2", []))

    # ============================================================
    # SLIDE 7 – AR RISK
    # ============================================================

    slide7 = prs.slides.add_slide(prs.slide_layouts[6])

    save_line_chart(ar["collection_rate_pct"],
                    "temp/ar_s3_rate.png",
                    "Collection Rate %")

    slide7.shapes.add_picture("temp/ar_s3_rate.png",
                              Inches(0.5), Inches(1),
                              width=Inches(6), height=Inches(3))

    add_insights(slide7, insights.get("ar", {}).get("slide3", []))

    # ============================================================
    # SLIDE 8 – EXECUTIVE SUMMARY
    # ============================================================

    slide8 = prs.slides.add_slide(prs.slide_layouts[6])

    add_insights(slide8, insights.get("executive_summary", []))

    output_path = "outputs/Finance_Control_Tower.pptx"
    prs.save(output_path)

    return output_path
