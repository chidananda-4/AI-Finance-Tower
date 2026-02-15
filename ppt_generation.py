# Commented out IPython magic to ensure Python compatibility.
import pandas as pd
import numpy as np
from openai import OpenAI
import os
import matplotlib.pyplot as plt
# %matplotlib inline
from pptx import Presentation
from pptx.util import Inches,Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from io import BytesIO
import io
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
import matplotlib.ticker as ticker
from pptx.enum.text import MSO_AUTO_SIZE,MSO_VERTICAL_ANCHOR,MSO_ANCHOR
import seaborn as sns
import json
from images_and_insights import *  # Your friend's visualization functions
from agents.summary_agent import SummaryAgent  # YOUR CrewAI Summary Agent
from datetime import datetime

def build_ap_view(excel_file):
    ap = excel_file["AP_Spend_Invoices"]
    vendor = excel_file["Vendor_Master"]
    cost = excel_file["CostCenter_Master"]
    budget = excel_file["Budget_Plan_Monthly"]
    kpi = excel_file["Finance_KPI_Weekly"]  # RESTORED

    ap['txn_date'] = pd.to_datetime(ap['txn_date'])
    ap['month'] = ap['txn_date'].dt.strftime('%Y-%m')
    ap['week'] = ap['txn_date'].dt.to_period('W')
    budget['month'] = pd.to_datetime(budget['month']).dt.to_period('M')

    kpi['week_start'] = pd.to_datetime(kpi['week_start'])
    kpi['maverick_pct'] = kpi['maverick_spend_inr'] / kpi['total_spend_inr'] * 100
    kpi['late_pct'] = kpi['late_payment_rate_ap'] * 100

    # Merge vendor & cost center
    ap = ap.merge(vendor, on=['vendor_id','payment_terms_days','preferred_vendor_flag'], how='left')
    ap = ap.merge(cost, on=['cost_center','business_unit','region'], how='left')

    # Return all three values
    return ap, budget, kpi  # RESTORED to 3 values

def build_ar_view(excel_file):
    ar_invoices = excel_file['AR_Invoices_Collections']
    #-------- AR SUMMARY------------------------------------------------------------
    ar_invoices['month_invoice'] = ar_invoices['invoice_date'].dt.strftime('%Y-%m')
    ar_invoices['month_paid'] = ar_invoices['paid_date'].dt.strftime('%Y-%m')
    ar_invoices['month_expected'] = ar_invoices['due_date'].dt.strftime('%Y-%m')
    ar_invoices['open_ar_flag'] = ar_invoices['payment_status'].apply(lambda x:1 if x=='Open' else 0)
    ar_invoices['paid_ar_flag'] = ar_invoices['payment_status'].apply(lambda x:1 if x!='Open' else 0)

    invoice_summary = ar_invoices.groupby(['month_invoice','segment','region']).agg(
        ar_invoiced_inr = ('invoice_amount_inr','sum'),
        open_ar_count = ('open_ar_flag','sum'),
        paid_ar_count = ('paid_ar_flag','sum'),
        late_payments = ('late_payment_flag','sum'),
        invoices = ('ar_invoice_id','count')).reset_index()

    paid_summary = ar_invoices[ar_invoices['payment_status']!='Open'].groupby(
        ['month_paid','segment','region']).agg(
        ar_collected_inr = ('paid_amount_inr','sum')).reset_index()

    expected_summary = ar_invoices.groupby(['month_expected','segment','region']).agg(
        ar_expected_inr = ('invoice_amount_inr','sum')).reset_index()

    ar_summary = invoice_summary.merge(paid_summary, 
        left_on = ['month_invoice','segment','region'],
        right_on = ['month_paid','segment','region'],how='left')

    ar_summary = ar_summary.merge(expected_summary, 
        left_on = ['month_invoice','segment','region'],
        right_on = ['month_expected','segment','region'],how='left')
    
    ar_summary = ar_summary[['month_invoice', 'segment', 'region','ar_invoiced_inr', 
                            'open_ar_count','paid_ar_count', 'late_payments', 
                            'invoices', 'ar_collected_inr']]
    ar_summary.rename(columns={'month_invoice':'month'},inplace=True)
    return ar_summary

def generate_insight_from_prompt(insight, n_images, client):
    """
    Generate insight using OpenAI based on number of images
    """
    max_bullets = 5

    full_prompt = f"""
    You are a senior business analyst.

    From the provided insight context below, generate executive-style insight lines.

    Rules:
    - No assumptions
    - No storytelling
    - No extra explanation
    - Each line must be concise.
    - No Grammetical mistakes
    - Output ONLY bullet-style lines separated by newline with {max_bullets} maximum bullet points which covers all the points in the insight.

    Insight Context:
    {insight}
    """

    response = client.chat.completions.create(
        model="gpt-4.1-mini",
        messages=[
            {"role": "system", "content": "Generate strictly factual business insights only from provided chart context."},
            {"role": "user", "content": full_prompt}
        ],
        temperature=0.0
    )

    return response.choices[0].message.content.strip()

# ---------------------------------------------------
def auto_font_size(text):
    length = len(text)

    if length < 200:
        return Pt(14)
    elif length < 400:
        return Pt(12)
    elif length < 700:
        return Pt(10)
    else:
        return Pt(9)

# ---------------------------------------------------
# DYNAMIC IMAGE + INSIGHT LAYOUT (NO OVERLAP) - From your friend
# ---------------------------------------------------
def get_dynamic_layout(n, slide_width, slide_height):

    margin_x = Inches(0.6)
    margin_top = Inches(1.4)
    margin_bottom = Inches(0.4)

    usable_w = slide_width - 2 * margin_x
    usable_h = slide_height - margin_top - margin_bottom

    coords = []
    insight_box = None

    # -----------------------------
    # CASE 1: ONE CHART
    # -----------------------------
    if n == 1:
        img_w = usable_w * 0.75
        img_h = usable_h * 0.65

        coords = [(margin_x + usable_w * 0.12, margin_top)]

        insight_box = (
            margin_x,
            margin_top + img_h + Inches(0.3),
            usable_w,
            Inches(1.5)
        )

    # -----------------------------
    # CASE 2: TWO CHARTS
    # -----------------------------
    elif n == 2:
        img_w = usable_w / 2 - Inches(0.25)
        img_h = usable_h * 0.6

        coords = [
            (margin_x, margin_top),
            (margin_x + img_w + Inches(0.4), margin_top)
        ]

        insight_box = (
            margin_x,
            margin_top + img_h + Inches(0.3),
            usable_w,
            Inches(1.5)
        )

    # -----------------------------
    # CASE 3: THREE CHARTS FIXED - From your friend
    # Layout: Chart1 | Chart2
    #         Insight | Chart3
    # -----------------------------
    elif n == 3:
        img_w = usable_w / 2 - Inches(0.25)
        img_h = usable_h * 0.42

        # Top row charts
        coords = [
            (margin_x, margin_top),                             # Chart 1
            (margin_x + img_w + Inches(0.4), margin_top)        # Chart 2
        ]

        # Bottom row split
        bottom_top = margin_top + img_h + Inches(0.35)

        insight_w = usable_w / 2 - Inches(0.25)
        insight_h = usable_h * 0.42

        # Chart3 goes to right side
        chart3_left = margin_x + insight_w + Inches(0.4)
        coords.append((chart3_left, bottom_top))

        # Insight box goes left side
        insight_box = (
            margin_x,
            bottom_top,
            insight_w,
            insight_h
        )

    # -----------------------------
    # CASE 4: FOUR CHARTS
    # -----------------------------
    else:
        img_w = usable_w / 2 - Inches(0.25)
        img_h = usable_h / 2 - Inches(0.25)

        coords = [
            (margin_x, margin_top),
            (margin_x + img_w + Inches(0.4), margin_top),
            (margin_x, margin_top + img_h + Inches(0.3)),
            (margin_x + img_w + Inches(0.4), margin_top + img_h + Inches(0.3))
        ]

        insight_box = (
            margin_x,
            margin_top + img_h * 2 + Inches(0.4),
            usable_w,
            Inches(1.3)
        )

    return coords, img_w, img_h, insight_box

# ---------------------------------------------------
# COVER PAGE - Option 4 with Quote
# ---------------------------------------------------
def add_cover_page(prs, latest_month):
    """
    Add cover page with quote and latest month
    """
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Main Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(1.2))
    tf = title_box.text_frame
    tf.text = "FINANCE CONTROL TOWER"
    tf.paragraphs[0].font.size = Pt(48)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(30, 74, 118)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(0.8))
    tf = sub_box.text_frame
    tf.text = "Autonomous Analytics Reporting System"
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.italic = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Report Period
    period_box = slide.shapes.add_textbox(Inches(1.5), Inches(5.5), Inches(12), Inches(0.6))
    tf = period_box.text_frame
    tf.text = f"Executive Report | {latest_month}"
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Footer
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12), Inches(0.3))
    tf = footer_box.text_frame
    tf.text = "CONFIDENTIAL - FOR INTERNAL USE ONLY"
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.color.rgb = RGBColor(150, 150, 150)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return slide

# ---------------------------------------------------
# THANK YOU PAGE - Matching Option 4 style
# ---------------------------------------------------
def add_thank_you_page(prs):
    """
    Add thank you page with quote and next steps
    """
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Thank You
    thank_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(12), Inches(1.2))
    tf = thank_box.text_frame
    tf.text = "THANK YOU"
    tf.paragraphs[0].font.size = Pt(54)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(30, 74, 118)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Appreciation
    app_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(12), Inches(0.8))
    tf = app_box.text_frame
    tf.text = "For your time and attention to this report"
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.italic = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Next Steps
    next_box = slide.shapes.add_textbox(Inches(1.5), Inches(6.0), Inches(12), Inches(0.6))
    tf = next_box.text_frame
    tf.text = "Let's turn insights into action"
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(30, 74, 118)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return slide

# ---------------------------------------------------
# YOUR Executive Summary Slide Creator - STORYTELLING FORMAT
# ---------------------------------------------------
def create_executive_summary_slide(prs, summary):
    """
    Storytelling narrative format executive summary
    Sections: The Big Picture | What's Working | What Needs Attention | Next Steps
    """
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
    
    # Create the slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # ==================== TITLE ====================
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(12), Inches(0.6))
    tf = title_box.text_frame
    tf.text = "EXECUTIVE SUMMARY"
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(30, 74, 118)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # ==================== THE BIG PICTURE ====================
    y_pos = 0.7
    section_title = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(12), Inches(0.3))
    tf = section_title.text_frame
    tf.text = "THE BIG PICTURE"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(30, 74, 118)
    
    y_pos += 0.3
    big_picture = summary.get("big_picture", "")
    if not big_picture:
        big_picture = "Q4 performance analysis completed. Review detailed slides for complete picture."
    
    bp_box = slide.shapes.add_textbox(Inches(0.6), Inches(y_pos), Inches(11.8), Inches(0.9))
    tf = bp_box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = big_picture
    p.font.size = Pt(11)
    p.font.italic = True
    p.line_spacing = 1.2
    p.space_after = Pt(0)
    
    # ==================== DIVIDER 1 ====================
    y_pos += 1.2
    line_left = Inches(0.5)
    line_top = Inches(y_pos)
    line_width = Inches(12)
    line_height = Inches(0.02)
    line = slide.shapes.add_shape(1, line_left, line_top, line_width, line_height)
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(200, 200, 200)
    line.line.color.rgb = RGBColor(200, 200, 200)
    
    y_pos += 0.25
    
    # ==================== TWO-COLUMN LAYOUT ====================
    col_width = Inches(5.6)
    col_left = Inches(0.6)
    col_right = Inches(6.4)
    
    # Left Column - WHAT'S WORKING
    left_header = slide.shapes.add_textbox(col_left, Inches(y_pos), col_width, Inches(0.3))
    tf = left_header.text_frame
    tf.text = "✓ WHAT'S WORKING"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(0, 100, 0)
    
    # Right Column - WHAT NEEDS ATTENTION
    right_header = slide.shapes.add_textbox(col_right, Inches(y_pos), col_width, Inches(0.3))
    tf = right_header.text_frame
    tf.text = "⚠ WHAT NEEDS ATTENTION"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(180, 0, 0)
    
    y_pos += 0.35
    
    # Left Column Content - Working
    working = summary.get("working", [])
    if not working:
        working = ["• Review detailed slides for positive developments"]
    
    left_box = slide.shapes.add_textbox(col_left, Inches(y_pos), col_width, Inches(1.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    
    for i, item in enumerate(working[:3]):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(10)
        p.space_after = Pt(6)
        p.line_spacing = 1.2
    
    # Right Column Content - Attention
    attention = summary.get("attention", [])
    if not attention:
        attention = ["• Review detailed slides for risk assessment"]
    
    right_box = slide.shapes.add_textbox(col_right, Inches(y_pos), col_width, Inches(1.5))
    tf = right_box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    
    for i, item in enumerate(attention[:3]):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(10)
        p.space_after = Pt(6)
        p.line_spacing = 1.2
    
    # ==================== DIVIDER 2 ====================
    y_pos += 1.8
    line_top = Inches(y_pos)
    line = slide.shapes.add_shape(1, line_left, line_top, line_width, line_height)
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(200, 200, 200)
    line.line.color.rgb = RGBColor(200, 200, 200)
    
    y_pos += 0.30
    
    # ==================== NEXT STEPS ====================
    next_header = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(12), Inches(0.3))
    tf = next_header.text_frame
    tf.text = "🎯 NEXT STEPS"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(30, 74, 118)
    
    y_pos += 0.35
    
    # Next Steps Content
    next_steps = summary.get("next_steps", [])
    if not next_steps:
        next_steps = [
            "THIS WEEK: Review budget variances. Finance Team.",
            "THIS MONTH: Optimize payment terms. Treasury.",
            "THIS QUARTER: Strategic planning. Leadership."
        ]
    
    next_box = slide.shapes.add_textbox(Inches(0.6), Inches(y_pos), Inches(11.8), Inches(1.2))
    tf = next_box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    
    for i, step in enumerate(next_steps[:4]):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = step
        p.font.size = Pt(10)
        p.space_after = Pt(6)
        p.line_spacing = 1.2
        if step.startswith("THIS WEEK"):
            p.font.bold = True
        else:
            p.font.bold = False
    
    # ==================== FOOTER ====================
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.2), Inches(12), Inches(0.2))
    tf = footer_box.text_frame
    tf.text = "Generated by Finance Control Tower AI | Data from latest month"
    tf.paragraphs[0].font.size = Pt(8)
    tf.paragraphs[0].font.color.rgb = RGBColor(150, 150, 150)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return slide

# ---------------------------------------------------
# Function to extract latest month from data
# ---------------------------------------------------
def get_latest_month(excel_file):
    """
    Extract the latest month from the data
    """
    try:
        # Try to get from AP data first
        if "AP_Spend_Invoices" in excel_file:
            ap_df = excel_file["AP_Spend_Invoices"]
            if 'txn_date' in ap_df.columns:
                ap_df['txn_date'] = pd.to_datetime(ap_df['txn_date'])
                latest = ap_df['txn_date'].max()
                return latest.strftime('%B %Y')
        
        # Fallback to AR data
        if "AR_Invoices_Collections" in excel_file:
            ar_df = excel_file["AR_Invoices_Collections"]
            if 'invoice_date' in ar_df.columns:
                ar_df['invoice_date'] = pd.to_datetime(ar_df['invoice_date'])
                latest = ar_df['invoice_date'].max()
                return latest.strftime('%B %Y')
    except:
        pass
    
    # Default fallback
    return datetime.now().strftime('%B %Y')

# ---------------------------------------------------
# MAIN PPT GENERATOR - WITH KPI RESTORED
# ---------------------------------------------------
def ppt_generator(excel_file, client, generate_summary=True):

    # ---- Build Data Views (RESTORED with kpi) ----
    ap, budget, kpi = build_ap_view(excel_file)  # Now returns 3 values
    ar_summary = build_ar_view(excel_file)

    # ---- Slide Data - USING YOUR FRIEND'S FUNCTIONS ----
    slides_data = [
        slide1_content(ap, ar_summary, budget, kpi, client),  # RESTORED kpi
        slide2_content(ap, ar_summary, budget, kpi, client),  # RESTORED kpi
        slide3_content(ap, ar_summary, budget, kpi, client),  # RESTORED kpi
        slide4_content(ap, ar_summary, budget, kpi, client),  # RESTORED kpi
        slide5_content(ap, ar_summary, budget, kpi, client),  # RESTORED kpi
        slide6_content(ap, ar_summary, budget, kpi, client),  # RESTORED kpi
        slide7_content(ap, ar_summary, budget, kpi, client),  # RESTORED kpi
        slide8_content(ap, ar_summary, budget, kpi, client),  # RESTORED kpi
        slide9_content(ap, ar_summary, budget, kpi, client),  # RESTORED kpi
        slide10_content(ap, ar_summary, budget, kpi, client), # RESTORED kpi
        slide11_content(ap, ar_summary, budget, kpi, client)  # RESTORED kpi
    ]

    # ---- Slide Titles ----
    title_dict = {
        1: "Monthly Spend Overview",
        2: "Spend Across Regions by Business Unit",
        3: "Spend Across Categories and Sub-Categories",
        4: "Budget Variance and Spend Utilization Summary",
        5: "Spend Analysis by Vendor Payment Terms & Regional Risk",
        6: "Maverick Spend Analysis",
        7: "Late Payment Insights",
        8: "Monthly Invoiced Amount Overview",
        9: "Invoiced Amount Trend by Region",
        10: "Monthly Collections Overview",
        11: "Collections Trend by Region"
    }

    # ---- Create Presentation ----
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    print("\n📊 Generating PPT Slides...\n")

    # ================================
    # ADD COVER PAGE (FIRST SLIDE)
    # ================================
    latest_month = get_latest_month(excel_file)
    print("📌 Adding Cover Page...")
    add_cover_page(prs, latest_month)

    # ================================
    # COLLECT INSIGHTS FOR CREWAI
    # ================================
    combined_insights = ""
    slides_content = []
    
    for i in range(1, 12):
        images, raw_text = slides_data[i-1]
        
        # Add to combined insights for CrewAI
        combined_insights += f"\n--- {title_dict[i]} ---\n"
        combined_insights += raw_text
        combined_insights += "\n"
        
        # Store slide content
        slides_content.append({
            "title": title_dict[i],
            "images": images,
            "insight": raw_text,
            "num_images": len(images)
        })
    
    # ================================
    # CREWAI Executive Summary
    # ================================
    executive_summary = None
    if generate_summary:
        print("\n🤖 CrewAI: Generating Executive Summary...")
        summary_agent = SummaryAgent(client)
        executive_summary = summary_agent.generate_executive_summary(combined_insights)
        
        # Add executive summary as next slide (after cover)
        print("📌 Adding Executive Summary Slide...")
        create_executive_summary_slide(prs, executive_summary)
    
    # ================================
    # ADD OTHER VISUALIZATION SLIDES
    # ================================
    for idx, slide_content in enumerate(slides_content, 1):
        slide_num = idx + 2 if generate_summary else idx + 1  # +2 for cover + exec summary
        print(f"\nGenerating Slide {slide_num}: {slide_content['title']}")
        
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.2),
            prs.slide_width - Inches(1), Inches(0.7)
        )
        tf = title_box.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = slide_content['title']
        p.font.size = Pt(20)
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT

        # Dynamic layout
        n = slide_content['num_images']
        coords, IMG_W, IMG_H, insight_box = get_dynamic_layout(
            n, prs.slide_width, prs.slide_height
        )

        # Add charts
        for img, (left, top) in zip(slide_content['images'], coords):
            img.seek(0)
            slide.shapes.add_picture(img, left, top, width=IMG_W, height=IMG_H)

        # Generate insights
        insight_text = generate_insight_from_prompt(slide_content['insight'], n, client)
        font_size = auto_font_size(insight_text)

        # Add insight box
        if insight_box:
            left, top, width, height = insight_box
            box = slide.shapes.add_textbox(left, top, width, height)
            fill = box.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(242, 242, 242)

            tf2 = box.text_frame
            tf2.word_wrap = True
            tf2.clear()
            tf2.vertical_anchor = MSO_ANCHOR.TOP

            para = tf2.paragraphs[0]
            para.clear()

            run1 = para.add_run()
            run1.text = "Insights:\n"
            run1.font.bold = True
            run1.font.size = font_size

            run2 = para.add_run()
            run2.text = insight_text
            run2.font.bold = False
            run2.font.size = font_size

            para.alignment = PP_ALIGN.LEFT

        print(f"✅ Slide {slide_num} completed")

    # ================================
    # ADD THANK YOU PAGE (LAST SLIDE)
    # ================================
    print("\n📌 Adding Thank You Page...")
    add_thank_you_page(prs)

    # Format insights for return
    insights = ""
    if executive_summary:
        insights += "EXECUTIVE SUMMARY (CrewAI Generated)\n"
        insights += "="*50 + "\n\n"
        insights += executive_summary.get("formatted_text", "")
        insights += "\n\n" + "="*50 + "\n\n"
    
    insights += "DETAILED SLIDE INSIGHTS\n"
    insights += "="*50 + "\n\n"

    total_slides = len(prs.slides)
    print("\n" + "="*60)
    print("✅ PPT Created Successfully!")
    print("="*60)
    print(f"   • Cover Page: Slide 1")
    if generate_summary:
        print(f"   • Executive Summary: Slide 2")
        print(f"   • Visualization Slides: Slides 3-{total_slides-1}")
    else:
        print(f"   • Visualization Slides: Slides 2-{total_slides-1}")
    print(f"   • Thank You Page: Slide {total_slides}")
    print(f"   • Total Slides: {total_slides}")

    return prs, insights, executive_summary